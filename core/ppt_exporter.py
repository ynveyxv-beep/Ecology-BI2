from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import tempfile


class PPTExporter:
    """Экспорт дашборда в PowerPoint"""
    
    @staticmethod
    def export(scene, path, title="Dashboard Report"):
        """Экспортирует сцену в PowerPoint"""
        try:
            prs = Presentation()
            
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(0.5), Inches(8), Inches(1)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(36)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            from PySide6.QtGui import QPixmap, QPainter
            
            scene_rect = scene.sceneRect()
            pixmap = QPixmap(int(scene_rect.width()), int(scene_rect.height()))
            pixmap.fill('white')
            
            painter = QPainter(pixmap)
            scene.render(painter)
            painter.end()
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            temp_path = temp_file.name
            temp_file.close()
            pixmap.save(temp_path, "PNG")
            
            slide.shapes.add_picture(
                temp_path,
                Inches(1), Inches(1.5),
                width=Inches(8),
                height=Inches(5.5)
            )
            
            prs.save(path)
            
            os.unlink(temp_path)
            
            print(f"✅ Exported to PowerPoint: {path}")
            return True
            
        except Exception as e:
            print(f"❌ Error exporting to PowerPoint: {e}")
            import traceback
            traceback.print_exc()
            return False